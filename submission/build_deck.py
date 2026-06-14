"""Build SUBMISSION.pptx for the ApplyForge capstone — screenshots + rubric mapping."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
SHOTS = os.path.join(HERE, "screenshots")
CROP = os.path.join(HERE, "_cropped")
os.makedirs(CROP, exist_ok=True)

INDIGO = RGBColor(0x31, 0x2E, 0x81)
INDIGO2 = RGBColor(0x4F, 0x46, 0xE5)
INK = RGBColor(0x1F, 0x29, 0x37)
GRAY = RGBColor(0x6B, 0x72, 0x80)
GREEN = RGBColor(0x15, 0x80, 0x3D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF3, 0xF4, 0xF6)

EMU_PER_IN = 914400


def crop_bottom(name):
    """Trim trailing near-white whitespace from a screenshot so it embeds tightly."""
    src = os.path.join(SHOTS, name)
    img = Image.open(src).convert("RGB")
    w, h = img.size
    px = img.load()
    bg = px[2, 2]
    last = h - 1
    for y in range(h - 1, -1, -1):
        row_has_content = False
        for x in range(0, w, 7):
            r, g, b = px[x, y]
            if abs(r - bg[0]) + abs(g - bg[1]) + abs(b - bg[2]) > 24:
                row_has_content = True
                break
        if row_has_content:
            last = min(h - 1, y + 30)
            break
    out = os.path.join(CROP, name)
    img.crop((0, 0, w, last + 1)).save(out)
    return out


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=None):
    """runs: list of paragraphs; each paragraph is list of (txt, size, bold, color)."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if space is not None:
            p.space_after = Pt(space)
        for (txt, size, bold, color) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
            r.font.name = "Calibri"
    return tb


def bullets(s, x, y, w, h, items, size=16, color=INK, gap=8):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        if isinstance(it, tuple):
            label, rest = it
            r = p.add_run(); r.text = "• " + label; r.font.bold = True
            r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = "Calibri"
            r2 = p.add_run(); r2.text = rest; r2.font.size = Pt(size)
            r2.font.color.rgb = color; r2.font.name = "Calibri"
        else:
            r = p.add_run(); r.text = "• " + it
            r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = "Calibri"
    return tb


def header(s, title, sub=None):
    rect(s, 0, 0, SW, Inches(1.15), INDIGO)
    text(s, Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.9),
         [[(title, 30, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
    if sub:
        text(s, Inches(0.52), Inches(1.25), Inches(12.3), Inches(0.5),
             [[(sub, 15, False, GRAY)]])


def place_image(s, path, x, y, max_w, max_h, border=True):
    img = Image.open(path); iw, ih = img.size
    scale = min(max_w / iw, max_h / ih)
    w = int(iw * scale); h = int(ih * scale)
    pic = s.shapes.add_picture(path, x + (max_w - w) // 2, y, width=w, height=h)
    if border:
        pic.line.color.rgb = RGBColor(0xD1, 0xD5, 0xDB); pic.line.width = Pt(1)
    return pic


def shot_slide(title, sub, shot, caption_bullets):
    s = slide()
    header(s, title, sub)
    img_path = crop_bottom(shot)
    # image on left ~7.6in wide, bullets on right
    place_image(s, img_path, Inches(0.4), Inches(1.5), Inches(7.7), Inches(5.6))
    text(s, Inches(8.4), Inches(1.6), Inches(4.5), Inches(0.5),
         [[("What this shows", 17, True, INDIGO2)]])
    bullets(s, Inches(8.4), Inches(2.2), Inches(4.6), Inches(5.0), caption_bullets, size=14, gap=10)
    return s


# ---------------- S1 Title ----------------
s = slide()
rect(s, 0, 0, SW, SH, INDIGO)
rect(s, 0, Inches(3.05), SW, Inches(0.06), INDIGO2)
text(s, Inches(1), Inches(2.0), Inches(11.3), Inches(1.1),
     [[("ApplyForge", 54, True, WHITE)]], align=PP_ALIGN.CENTER)
text(s, Inches(1), Inches(3.2), Inches(11.3), Inches(0.7),
     [[("An honest AI job-search copilot", 24, False, RGBColor(0xC7,0xD2,0xFE))]],
     align=PP_ALIGN.CENTER)
text(s, Inches(1), Inches(4.5), Inches(11.3), Inches(1.4),
     [[("CSE 552 — Capstone Project", 18, True, WHITE)],
      [("Full-stack AI web app · Next.js + FastAPI + PostgreSQL + Claude · Docker Compose", 14, False, RGBColor(0xC7,0xD2,0xFE))],
      [("Tray Branch · June 2026", 14, False, RGBColor(0xC7,0xD2,0xFE))]],
     align=PP_ALIGN.CENTER, space=6)

# ---------------- S2 Problem & Solution ----------------
s = slide()
header(s, "The problem & the solution", "What you built — the problem, the user, the solution")
text(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.5),
     [[("The user: ", 18, True, INDIGO2), ("a software engineer job-hunting, drowning in applications.", 18, False, INK)]])
bullets(s, Inches(0.5), Inches(2.2), Inches(12.3), Inches(2.0), [
    ("The problem: ", "tools like RoboApply/ApplyBlast auto-blast hundreds of applications — recruiters ignore the spam, and it violates job-site terms."),
    ("The insight (from research): ", "the “an AI robot rejects 75% of resumes” claim is a myth (traces to a defunct 2012 vendor). The real levers are honest tailoring, timing, and direct recruiter contact."),
    ("The solution: ", "ApplyForge tells you your REAL match for a role, tailors your resume using only what's actually true about you, and helps you reach a real human."),
], size=16, gap=12)
rect(s, Inches(0.5), Inches(5.7), Inches(12.3), Inches(1.1), LIGHT)
text(s, Inches(0.8), Inches(5.85), Inches(11.7), Inches(0.9),
     [[("Why it's “AI-powered”: ", 16, True, INDIGO2),
       ("six Claude endpoints inject the user's own profile + job + pipeline as context — match scoring, "
        "resume tailoring, cover letters, recruiter outreach, and a tool-using agent. Not a generic chatbot.", 15, False, INK)]],
     anchor=MSO_ANCHOR.MIDDLE)

# ---------------- S3 Architecture ----------------
s = slide()
header(s, "Architecture", "Three-tier stack, containerized with Docker Compose")
mono = [
    [("Browser  ──HTTP──▶  FastAPI  ──SQLAlchemy──▶  PostgreSQL 16", 16, True, INK)],
    [("  :3000              :8000                       :5432", 12, False, GRAY)],
    [("", 8, False, INK)],
    [("Next.js 14 (App Router, TS, Tailwind)   →   one typed fetch layer  app/lib/api.ts", 13, False, INK)],
    [("FastAPI + Pydantic + SQLAlchemy         →   every Claude call isolated in  ai.py", 13, False, INK)],
    [("Claude API: Sonnet 4.6 (reasoning) + Haiku 4.5 (transforms)", 13, False, INK)],
    [("docker compose up --build  →  db (healthcheck) → backend (auto-seed) → frontend", 13, True, INDIGO2)],
]
rect(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(2.6), LIGHT)
text(s, Inches(0.8), Inches(1.8), Inches(11.8), Inches(2.3), mono, space=4)
text(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.4), [[("Separation of concerns", 17, True, INDIGO2)]])
bullets(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.0), [
    "Every Claude call lives in backend/ai.py; every network call in frontend/app/lib/api.ts",
    "Job feed isolated in services/job_feed.py; file parsing in services/resume_text.py",
    "Data persists in a Docker volume (pgdata); seed runs only when the DB is empty",
], size=15, gap=9)

# ---------------- S4 Data model ----------------
s = slide()
header(s, "Data model", "7 SQLAlchemy models, PostgreSQL — well beyond the 2-model minimum")
diagram = [
    [("Profile (1) ──< Application (many) >── Job (many)", 18, True, INK)],
    [("                    │", 14, False, GRAY)],
    [("                    ├──< GeneratedArtifact   (match reports, cover letters, tailored resumes)", 13, False, INK)],
    [("                    └──< ApplicationEvent     (pipeline audit log)", 13, False, INK)],
    [("Profile (1) ──< Outreach (many) >── Recruiter (many)", 18, True, INK)],
]
rect(s, Inches(0.5), Inches(1.7), Inches(12.3), Inches(2.5), LIGHT)
text(s, Inches(0.8), Inches(2.0), Inches(11.8), Inches(2.0), diagram, space=8)
bullets(s, Inches(0.5), Inches(4.6), Inches(12.3), Inches(2.4), [
    ("Application ", "is the core CRUD resource: create (track a job), read, update (advance/reject status), delete."),
    ("Relationships: ", "one-to-many Profile→Application→Job, Recruiter→Outreach, plus artifacts & an append-only event log."),
    ("Pipeline state machine: ", "SAVED → APPLIED → INTERVIEWING → OFFER, with REJECTED as a terminal branch (forward-only, validated server-side)."),
], size=15, gap=11)

# ---------------- Screenshot slides ----------------
shot_slide("Dashboard — CRUD pipeline", "frontend page 1 of 6 · shared nav · responsive",
           "01-dashboard.png", [
    ("Stat cards", " from a live /stats aggregate: 5 applications, 71% avg match, outreach, jobs."),
    ("Pipeline board", " — every tracked application grouped by its lifecycle stage."),
    ("Color-coded match badges", " (green = strong, red = stretch) from the AI analysis."),
    ("Shared top nav", " across all 6 pages; fully responsive Tailwind layout."),
])

shot_slide("Jobs — real openings from public ATS APIs", "page 2 · external API integration",
           "03-jobs.png", [
    ("Live job feed", " pulled from Greenhouse / Lever / Ashby public APIs — no scraping, no keys, no ToS risk."),
    ("Ranked", " by a quick keyword match against the user's profile."),
    ("“Track”", " creates an Application (the CRUD create path)."),
    ("Paste-a-JD", " fallback + “Sync live jobs” button.")
])

shot_slide("Application detail — the AI hero", "page 3 · structured output + generation, both using app data",
           "04-application-detail.png", [
    ("Match analysis", " (Claude Sonnet, structured JSON): honest 84% score, matched vs MISSING keywords, strengths & gaps."),
    ("Cover letter", " generated from the user's real resume — note it cites actual projects & metrics, never invented."),
    ("Truthful tailoring", " guardrail enforced in the prompt layer."),
    ("Pipeline controls", " + append-only event timeline (audit log)."),
])

shot_slide("Profile — AI resume extraction", "page 4 · structured extraction from unstructured input",
           "02-profile.png", [
    ("Upload a PDF/DOCX/TXT resume", " → Claude extracts a structured profile (skills, experience, headline)."),
    ("This is the source of truth", " every other AI feature is grounded in."),
    ("Full CRUD edit", " of profile fields, skills, and links."),
    ("Loading / error / saving states", " throughout."),
])

shot_slide("Recruiters — AI outreach generation", "page 5 · personalized content generation",
           "05-recruiters.png", [
    ("Seeded niche recruiters", " with specialty tags."),
    ("“Generate outreach”", " → Claude writes a 50–125 word note referencing a REAL accomplishment + the recruiter's niche, with an opt-out courtesy."),
    ("Sending is simulated", " — honest by design, no spam."),
    ("Outreach history", " persists per recruiter."),
])

shot_slide("Copilot — AI agent with tool use", "page 6 · multi-turn agent that reads your data",
           "06-copilot.png", [
    ("A real tool-using agent", " (Claude Sonnet): tools get_applications, get_application_detail, get_profile."),
    ("Grounded in live data", " — asks “which roles am I weakest for?” and reads actual match scores to answer."),
    ("Multi-turn", " conversation with suggested prompts."),
    ("Shows which tools were called", " under each reply."),
])

# ---------------- AI endpoints table ----------------
s = slide()
header(s, "The AI integration — 6 Claude endpoints", "All inject the user's app data; cost-routed Sonnet / Haiku")
rows = [
    ("Endpoint", "Claude technique", "Context injected"),
    ("POST /profile/import", "Structured extraction", "uploaded resume text"),
    ("POST /applications/{id}/analyze  ⭐", "Structured output (JSON)", "profile + job description"),
    ("POST /applications/{id}/cover-letter", "Generation (truthful)", "profile + job"),
    ("POST /applications/{id}/tailor", "Generation, structured", "resume bullets + job"),
    ("POST /outreach/generate", "Generation (personalized)", "profile + recruiter + job"),
    ("POST /ai/chat", "Tool use / agent (multi-turn)", "whole pipeline via 3 tools"),
]
tbl_shape = s.shapes.add_table(len(rows), 3, Inches(0.5), Inches(1.5), Inches(12.3), Inches(4.2))
table = tbl_shape.table
table.columns[0].width = Inches(5.0); table.columns[1].width = Inches(4.0); table.columns[2].width = Inches(3.3)
for c in range(3):
    cell = table.cell(0, c)
    cell.fill.solid(); cell.fill.fore_color.rgb = INDIGO
    p = cell.text_frame.paragraphs[0]; r = p.add_run(); r.text = rows[0][c]
    r.font.bold = True; r.font.color.rgb = WHITE; r.font.size = Pt(13)
for ri in range(1, len(rows)):
    for c in range(3):
        cell = table.cell(ri, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if ri % 2 else LIGHT
        p = cell.text_frame.paragraphs[0]; r = p.add_run(); r.text = rows[ri][c]
        r.font.size = Pt(12); r.font.color.rgb = INK
        if c == 0: r.font.bold = True
text(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.0),
     [[("Meets “at least one of” in spades: ", 14, True, INDIGO2),
       ("structured output (analyze), tool use (chat agent), AND multi-turn conversation — all three.", 14, False, INK)]])

# ---------------- Technical highlight ----------------
s = slide()
header(s, "Technical highlight", "The copilot is a real tool-use agent — not a scripted chatbot")
bullets(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(3.2), [
    "backend/ai.py:copilot_chat runs a genuine agent loop with the Claude Messages API.",
    "Claude decides which tool to call; the backend executes it against PostgreSQL and feeds the JSON result back; the loop repeats until Claude has enough to answer.",
    "Tools: get_applications, get_application_detail, get_profile — so answers are grounded in the user's live pipeline, never guessed.",
    "Verified live: “which application has the highest match?” → called get_applications → correctly answered Brightwave (88), edging Stripe (84).",
], size=16, gap=12)
rect(s, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.6), LIGHT)
text(s, Inches(0.8), Inches(5.4), Inches(11.8), Inches(1.3),
     [[("Honesty enforced in code: ", 16, True, INDIGO2),
       ("a shared TRUTHFUL_GUARDRAIL system prompt is injected into every generation call, and the resume "
        "tailor returns a per-bullet “grounded” flag. The integrity claim is enforced at the prompt layer, not just marketing.", 15, False, INK)]],
     anchor=MSO_ANCHOR.MIDDLE)

# ---------------- Requirements checklist ----------------
s = slide()
header(s, "Minimum requirements — all met", "Every checkbox from the spec")
left = [
    ("FRONTEND (Next.js)", True),
    ("6 pages with shared nav (need 4+)", False),
    ("Full CRUD UI for Applications", False),
    ("AI features with working UI", False),
    ("Loading & error states", False),
    ("Responsive design + Dockerfile", False),
    ("BACKEND (FastAPI)", True),
    ("~20 endpoints (need 6+)", False),
    ("7 models w/ relationships (need 2+)", False),
    ("6 AI endpoints (need 2+)", False),
    ("CORS configured + Dockerfile", False),
]
right = [
    ("DATABASE", True),
    ("PostgreSQL 16 via SQLAlchemy", False),
    ("Persists via Docker volume (pgdata)", False),
    ("DEPLOYMENT", True),
    ("docker-compose.yml wires 3 services", False),
    ("Starts with docker compose up --build", False),
    (".env gitignored (.dockerignore strips it)", False),
    ("AI INTEGRATION", True),
    ("Claude API (Sonnet 4.6 + Haiku 4.5)", False),
    ("AI has user data via context injection", False),
    ("Structured output + tool use + multi-turn", False),
]
def checklist(s, items, x):
    tb = s.shapes.add_textbox(x, Inches(1.5), Inches(6.2), Inches(5.6)); tf = tb.text_frame; tf.word_wrap=True
    for i, (txt, hdr) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph(); p.space_after = Pt(6)
        if hdr:
            r = p.add_run(); r.text = txt; r.font.bold = True; r.font.size = Pt(14); r.font.color.rgb = INDIGO2
        else:
            r = p.add_run(); r.text = "✓ "; r.font.bold = True; r.font.size = Pt(14); r.font.color.rgb = GREEN
            r2 = p.add_run(); r2.text = txt; r2.font.size = Pt(13); r2.font.color.rgb = INK
checklist(s, left, Inches(0.5))
checklist(s, right, Inches(6.9))

# ---------------- Rubric mapping ----------------
s = slide()
header(s, "Grading rubric — where each point is earned")
rows = [
    ("Criteria", "Pts", "Evidence in ApplyForge"),
    ("Core CRUD works in production", "15", "Applications track/advance/reject/delete + Profile edit, live via docker compose"),
    ("AI feature meaningful & uses app data", "20", "6 Claude endpoints inject profile+job+pipeline (analyze, cover letter, tailor, outreach, agent)"),
    ("All major flows work without errors", "10", "Smoke-tested live; 14 backend tests pass; AI calls verified end-to-end"),
    ("Frontend: Next.js, multi-page, responsive", "10", "6 pages, shared nav, Tailwind responsive, loading/error states"),
    ("Backend: FastAPI, 2+ models, organized", "10", "7 models, ~20 endpoints, ai.py / services / routes separation"),
    ("Database: Postgres, persists on restart", "5", "PostgreSQL 16 + SQLAlchemy, pgdata volume, seed-only-if-empty"),
    ("AI: Claude integrated thoughtfully", "10", "Sonnet+Haiku cost routing, structured output, tool-use agent, truthful guardrail"),
    ("Docker: Dockerfiles + compose, clean start", "5", "3 services, healthcheck gate, docker compose up --build"),
]
tbl_shape = s.shapes.add_table(len(rows), 3, Inches(0.4), Inches(1.35), Inches(12.5), Inches(5.4))
table = tbl_shape.table
table.columns[0].width = Inches(4.4); table.columns[1].width = Inches(0.8); table.columns[2].width = Inches(7.3)
for c in range(3):
    cell = table.cell(0, c); cell.fill.solid(); cell.fill.fore_color.rgb = INDIGO
    p = cell.text_frame.paragraphs[0]; r = p.add_run(); r.text = rows[0][c]
    r.font.bold = True; r.font.color.rgb = WHITE; r.font.size = Pt(12)
for ri in range(1, len(rows)):
    for c in range(3):
        cell = table.cell(ri, c); cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if ri % 2 else LIGHT
        p = cell.text_frame.paragraphs[0]; r = p.add_run(); r.text = rows[ri][c]
        r.font.size = Pt(11); r.font.color.rgb = INK
        if c == 1: r.font.bold = True; r.font.color.rgb = GREEN; p.alignment = PP_ALIGN.CENTER

# ---------------- What's next + run ----------------
s = slide()
header(s, "What I'd do next  ·  How to run")
text(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.4), [[("If I had two more weeks", 18, True, INDIGO2)]])
bullets(s, Inches(0.5), Inches(2.05), Inches(12.3), Inches(2.0), [
    "A Chrome extension that maps my profile onto a job form's fields via a /autofill/map Claude endpoint — the SAFE version of auto-apply (review-and-submit, your data, your click).",
    "Real recruiter data through a licensed API (Apollo / People Data Labs) with proper consent + opt-out handling.",
    "Multi-user auth + streaming AI responses.",
], size=15, gap=10)
rect(s, Inches(0.5), Inches(4.6), Inches(12.3), Inches(2.2), LIGHT)
text(s, Inches(0.8), Inches(4.8), Inches(11.8), Inches(2.0), [
    [("Run it (the graded path)", 16, True, INDIGO2)],
    [("1.  Put ANTHROPIC_API_KEY in backend/.env", 14, False, INK)],
    [("2.  docker compose up --build", 14, True, INK)],
    [("3.  Open http://localhost:3000  (API docs at :8000/docs)", 14, False, INK)],
    [("Auto-seeds a full demo on first run; data persists across restarts via the pgdata volume.", 13, False, GRAY)],
], space=6)

prs.save(os.path.join(HERE, "ApplyForge_Capstone_Submission.pptx"))
print("Saved ApplyForge_Capstone_Submission.pptx with", len(prs.slides._sldIdLst), "slides")
